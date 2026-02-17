import streamlit as st
import os
import tempfile
import numpy as np
import cv2
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont
import hashlib
import json
import math
import re
import zipfile
from io import BytesIO

st.set_page_config(page_title="PPTX Video Factory", layout="wide")

LAYOUTS = {
    "Original PPTX Size": None,
    "Landscape (16:9) 1920x1080": {"w": 1920, "h": 1080},
    "Portrait (9:16) 1080x1920": {"w": 1080, "h": 1920},
    "Square (1:1) 1080x1080": {"w": 1080, "h": 1080},
}

HORIZONTAL_ALIGN_MAP = {
    PP_ALIGN.LEFT: 'left',
    PP_ALIGN.CENTER: 'center',
    PP_ALIGN.RIGHT: 'right',
    PP_ALIGN.JUSTIFY: 'justify',
    PP_ALIGN.DISTRIBUTE: 'justify',
    None: 'left'
}

VERTICAL_ALIGN_MAP = {
    MSO_ANCHOR.TOP: 'top',
    MSO_ANCHOR.MIDDLE: 'middle',
    MSO_ANCHOR.BOTTOM: 'bottom',
    None: 'top'
}

def find_font():
    candidates = [
        "poppins.ttf", "Poppins.ttf", "Poppins-Bold.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for f in candidates:
        if os.path.exists(f):
            try:
                ImageFont.truetype(f, 40)
                return f
            except:
                continue
    return None

def get_images_from_folder(folder_path):
    """Get all images from a folder."""
    folder = Path(folder_path)
    if not folder.exists():
        return {}
    
    extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff', '.tif']
    images = {}
    
    for ext in extensions:
        for f in folder.glob(f"*{ext}"):
            images[f.name] = str(f)
    
    return images

def extract_media_from_pptx(pptx_path, extract_dir):
    """Extract media folder from PPTX to directory."""
    media_dir = Path(extract_dir) / "ppt" / "media"
    media_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for name in zf.namelist():
                if name.startswith('ppt/media/'):
                    zf.extract(name, extract_dir)
        return str(media_dir)
    except Exception as e:
        st.error(f"Error extracting media: {e}")
        return None

def find_match(query_name, available_images):
    """Simple name matching."""
    if not available_images or not query_name:
        return None
    
    query_clean = query_name.lower().replace(' ', '').replace('_', '').replace('-', '')
    
    for img_name, img_path in available_images.items():
        img_clean = img_name.lower().replace(' ', '').replace('_', '').replace('-', '')
        
        # Exact match
        if query_clean == img_clean:
            return img_path
        
        # Contains
        if query_clean in img_clean or img_clean in query_clean:
            return img_path
    
    return None

def get_shape_identifier(shape):
    """Get identifier for matching."""
    name = shape.name
    
    alt_text = None
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            alt_text = shape._pic.nvPicPr.cNvPr.get('descr')
            if alt_text:
                alt_text = Path(alt_text).stem
        except:
            pass
    
    return {
        'name': name,
        'alt_text': alt_text,
    }

def get_suggested_filename(shape_name, alt_text):
    """Generate suggested filename for missing image."""
    if alt_text:
        base = Path(alt_text).stem
    else:
        base = shape_name
    
    base = re.sub(r'[^\w\s-]', '', base).strip()
    base = base.replace(' ', '_').replace('-', '_').lower()
    base = re.sub(r'_+', '_', base)
    
    return f"{base}.png" if base else "image.png"

def extract_text_frame_properties(text_frame, scale=1.0):
    """Extract text formatting."""
    if not text_frame or not text_frame.paragraphs:
        return None
    
    margin_left = int(text_frame.margin_left * 96 / 914400 * scale) if text_frame.margin_left else 0
    margin_right = int(text_frame.margin_right * 96 / 914400 * scale) if text_frame.margin_right else 0
    margin_top = int(text_frame.margin_top * 96 / 914400 * scale) if text_frame.margin_top else 0
    margin_bottom = int(text_frame.margin_bottom * 96 / 914400 * scale) if text_frame.margin_bottom else 0
    
    vertical_anchor = text_frame.vertical_anchor
    vertical_align = VERTICAL_ALIGN_MAP.get(vertical_anchor, 'top')
    
    paragraphs = []
    for para in text_frame.paragraphs:
        if not para.runs and not para.text.strip():
            continue
            
        para_align = para.alignment
        horizontal_align = HORIZONTAL_ALIGN_MAP.get(para_align, 'left')
        
        run = para.runs[0] if para.runs else None
        if run and run.font:
            font_size = int(run.font.size.pt * scale) if run.font.size else int(24 * scale)
            bold = bool(run.font.bold)
            italic = bool(run.font.italic)
            font_name = run.font.name or "Arial"
            
            color = "#000000"
            try:
                if run.font.color and run.font.color.rgb:
                    color = f"#{run.font.color.rgb}"
            except:
                pass
        else:
            font_size = int(24 * scale)
            bold = italic = False
            font_name = "Arial"
            color = "#000000"
        
        paragraphs.append({
            'text': para.text,
            'align': horizontal_align,
            'font_size': max(font_size, 8),
            'font_name': font_name,
            'bold': bold,
            'italic': italic,
            'color': color,
            'line_spacing': 1.2
        })
    
    return {
        'paragraphs': paragraphs,
        'vertical_align': vertical_align,
        'margins': {
            'left': margin_left,
            'right': margin_right,
            'top': margin_top,
            'bottom': margin_bottom
        }
    }

def extract_background(slide):
    try:
        fill = slide.background.fill
        if fill.type is not None:
            if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                rgb = fill.fore_color.rgb
                return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
    except:
        pass
    return (255, 255, 255)

def harvest_ppt(pptx_path, target_dims=None, asset_sources=None, manual_mappings=None):
    try:
        prs = Presentation(pptx_path)
        if not prs.slides:
            return None
        
        slide = prs.slides[0]
        emu_to_px = 96 / 914400
        
        orig_w = int(prs.slide_width * emu_to_px)
        orig_h = int(prs.slide_height * emu_to_px)
        
        if target_dims:
            target_w, target_h = target_dims['w'], target_dims['h']
            scale = max(target_w / orig_w, target_h / orig_h)
            scaled_w = int(orig_w * scale)
            scaled_h = int(orig_h * scale)
            offset_x = (target_w - scaled_w) // 2
            offset_y = (target_h - scaled_h) // 2
            canvas_w, canvas_h = target_w, target_h
        else:
            canvas_w, canvas_h = orig_w, orig_h
            scale = 1.0
            offset_x, offset_y = 0, 0
        
        bg_color = extract_background(slide)
        
        # Build auto-match pool: folder -> extracted -> uploaded
        auto_match_pool = {}
        
        if asset_sources:
            if asset_sources.get('folder'):
                auto_match_pool.update(get_images_from_folder(asset_sources['folder']))
            if asset_sources.get('extracted_media'):
                auto_match_pool.update(get_images_from_folder(asset_sources['extracted_media']))
            if asset_sources.get('uploaded'):
                auto_match_pool.update(asset_sources['uploaded'])
        
        config = {
            "canvas": {"w": canvas_w, "h": canvas_h},
            "original": {"w": orig_w, "h": orig_h},
            "scale": scale,
            "offset": {"x": offset_x, "y": offset_y},
            "pptx_background": bg_color,
            "asset_sources": asset_sources,
            "auto_match_pool": list(auto_match_pool.keys()),
            "elements": []
        }
        
        def process_shape(shape, parent_x=0, parent_y=0):
            identifiers = get_shape_identifier(shape)
            original_id = identifiers['name']
            display_name = original_id
            
            x = int((shape.left * emu_to_px * scale) + offset_x + parent_x)
            y = int((shape.top * emu_to_px * scale) + offset_y + parent_y)
            w = int(shape.width * emu_to_px * scale)
            h = int(shape.height * emu_to_px * scale)
            
            # Handle groups
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                children = []
                for child in shape.shapes:
                    child_offset_x = x - offset_x
                    child_offset_y = y - offset_y
                    child_result = process_shape(child, parent_x=child_offset_x, parent_y=child_offset_y)
                    
                    if isinstance(child_result, list):
                        children.extend(child_result)
                    elif child_result:
                        children.append(child_result)
                return children if children else None
            
            el = {
                "id": original_id,
                "name": display_name,
                "alt_text": identifiers.get('alt_text'),
                "x": x, "y": y, "w": w, "h": h,
                "rotation": getattr(shape, 'rotation', 0) or 0,
                "z_order": getattr(shape, 'z_order', 0) or 0,
            }
            
            # TEXT - Always extract
            if shape.has_text_frame and shape.text.strip():
                text_props = extract_text_frame_properties(shape.text_frame, scale)
                if text_props and text_props['paragraphs']:
                    el.update({
                        "type": "text",
                        "text_props": text_props,
                        "text_default": shape.text_frame.text,
                    })
                    return el
            
            # PICTURE
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                el["type"] = "image"
                
                if manual_mappings and original_id in manual_mappings:
                    el["image_path"] = manual_mappings[original_id]
                    el["match_source"] = "manual"
                else:
                    matched = find_match(display_name, auto_match_pool)
                    
                    if not matched and identifiers.get('alt_text'):
                        matched = find_match(identifiers['alt_text'], auto_match_pool)
                    
                    if matched:
                        el["image_path"] = matched
                        if asset_sources.get('folder') and matched in get_images_from_folder(asset_sources['folder']).values():
                            el["match_source"] = "folder"
                        elif asset_sources.get('extracted_media') and matched in get_images_from_folder(asset_sources['extracted_media']).values():
                            el["match_source"] = "extracted"
                        else:
                            el["match_source"] = "uploaded"
                    else:
                        el["suggested_filename"] = get_suggested_filename(display_name, identifiers.get('alt_text'))
                
                return el
            
            # SHAPE
            el["type"] = "shape"
            
            if manual_mappings and original_id in manual_mappings:
                el["type"] = "image"
                el["image_path"] = manual_mappings[original_id]
                el["match_source"] = "manual"
                return el
            
            matched = find_match(display_name, auto_match_pool)
            if not matched and identifiers.get('alt_text'):
                matched = find_match(identifiers['alt_text'], auto_match_pool)
            
            if matched:
                el["type"] = "image"
                el["image_path"] = matched
                if asset_sources.get('folder') and matched in get_images_from_folder(asset_sources['folder']).values():
                    el["match_source"] = "folder"
                elif asset_sources.get('extracted_media') and matched in get_images_from_folder(asset_sources['extracted_media']).values():
                    el["match_source"] = "extracted"
                else:
                    el["match_source"] = "uploaded"
                return el
            
            # Shape styling
            try:
                if hasattr(shape, 'fill') and shape.fill.type == 1:
                    if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                        rgb = shape.fill.fore_color.rgb
                        el["fill_color"] = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
            except:
                pass
            
            try:
                if shape.has_line and shape.line.color.rgb:
                    rgb = shape.line.color.rgb
                    el["line_color"] = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
                    el["line_width"] = int(shape.line.width.pt * scale) if shape.line.width else 1
            except:
                pass
            
            return el
        
        for shape in slide.shapes:
            result = process_shape(shape)
            if isinstance(result, list):
                config["elements"].extend(result)
            elif result:
                config["elements"].append(result)
        
        config["elements"] = [el for el in config["elements"] if isinstance(el, dict)]
        config["elements"].sort(key=lambda x: x.get('z_order', 0))
        
        return config
        
    except Exception as e:
        st.error(f"Extraction failed: {e}")
        import traceback
        st.code(traceback.format_exc())
        return None

_font_cache = {}

def get_font_cached(font_path, size, bold=False, italic=False):
    cache_key = (font_path, size, bold, italic)
    if cache_key in _font_cache:
        return _font_cache[cache_key]
    
    if font_path and os.path.exists(font_path):
        try:
            font = ImageFont.truetype(font_path, size)
        except:
            font = ImageFont.load_default()
    else:
        font = ImageFont.load_default()
    
    _font_cache[cache_key] = font
    return font

def render_text_paragraphs(text_props, box_w, box_h, font_path):
    img = Image.new("RGBA", (box_w, box_h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    
    paragraphs = text_props.get('paragraphs', [])
    if not paragraphs:
        return np.array(img)
    
    margins = text_props.get('margins', {'left': 0, 'right': 0, 'top': 0, 'bottom': 0})
    vertical_align = text_props.get('vertical_align', 'top')
    
    avail_w = box_w - margins['left'] - margins['right']
    avail_h = box_h - margins['top'] - margins['bottom']
    
    para_heights = []
    
    for para in paragraphs:
        font_size = para['font_size']
        line_height = int(font_size * para.get('line_spacing', 1.2))
        
        text = para['text']
        font = get_font_cached(font_path, font_size, para.get('bold', False), para.get('italic', False))
        
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = draw.textbbox((0, 0), test_line, font=font)
            text_width = bbox[2] - bbox[0]
            
            if text_width <= avail_w:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        if not lines:
            lines = [text]
        
        para['computed_lines'] = lines
        para_heights.append(len(lines) * line_height)
    
    total_content_height = sum(para_heights)
    
    if vertical_align == 'middle':
        start_y = margins['top'] + (avail_h - total_content_height) // 2
    elif vertical_align == 'bottom':
        start_y = margins['top'] + avail_h - total_content_height
    else:
        start_y = margins['top']
    
    current_y = start_y
    
    for i, para in enumerate(paragraphs):
        font_size = para['font_size']
        line_height = int(font_size * para.get('line_spacing', 1.2))
        align = para.get('align', 'left')
        color_hex = para.get('color', '#000000')
        
        color_hex = color_hex.lstrip('#')
        color = tuple(int(color_hex[j:j+2], 16) for j in (0, 2, 4)) + (255,)
        
        font = get_font_cached(font_path, font_size, para.get('bold', False), para.get('italic', False))
        lines = para.get('computed_lines', [para['text']])
        
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=font)
            text_width = bbox[2] - bbox[0]
            
            if align == 'center':
                x = margins['left'] + (avail_w - text_width) // 2
            elif align == 'right':
                x = margins['left'] + avail_w - text_width
            else:
                x = margins['left']
            
            draw.text((x, current_y - bbox[1]), line, fill=color, font=font)
            current_y += line_height
        
        current_y += int(font_size * 0.3)
    
    return np.array(img)

def render_frame(layout, user_data, font_path, bg_settings):
    w, h = layout['canvas']['w'], layout['canvas']['h']
    
    bg_type = bg_settings.get('type', 'pptx')
    bg_value = bg_settings.get('value')
    
    if bg_type == 'image' and isinstance(bg_value, Image.Image):
        frame = np.array(bg_value.resize((w, h), Image.Resampling.LANCZOS))
    elif bg_type == 'color' and isinstance(bg_value, tuple):
        frame = np.full((h, w, 3), bg_value, dtype=np.uint8)
    elif bg_type == 'pptx':
        frame = np.full((h, w, 3), layout['pptx_background'], dtype=np.uint8)
    else:
        frame = np.full((h, w, 3), (245, 245, 245), dtype=np.uint8)
    
    pil_img = Image.fromarray(frame)
    draw = ImageDraw.Draw(pil_img)
    
    for el in layout['elements']:
        x, y, ew, eh = el['x'], el['y'], el['w'], el['h']
        
        if x < -ew or y < -eh or x >= w or y >= h:
            continue
        
        if el['type'] == 'shape':
            fill = el.get('fill_color', (200, 200, 200))
            line = el.get('line_color', (100, 100, 100))
            line_w = el.get('line_width', 1)
            
            if not isinstance(fill, tuple):
                fill = (200, 200, 200)
            if not isinstance(line, tuple):
                line = (100, 100, 100)
            
            draw.rectangle([x, y, x+ew, y+eh], fill=fill, outline=line, width=max(line_w, 1))
        
        elif el['type'] == 'image':
            img_path = el.get('image_path')
            
            if img_path and Path(img_path).exists():
                try:
                    asset_img = Image.open(img_path).convert('RGBA')
                    asset_img = asset_img.resize((ew, eh), Image.Resampling.LANCZOS)
                    pil_img.paste(asset_img, (x, y), asset_img)
                except Exception as e:
                    draw.rectangle([x, y, x+ew, y+eh], fill=(255, 0, 0))
                    draw.text((x+5, y+5), "ERR", fill=(255, 255, 255))
            else:
                # MISSING IMAGE - Show placeholder
                draw.rectangle([x, y, x+ew, y+eh], fill=(220, 220, 220), outline=(255, 100, 100), width=2)
                
                suggestion = el.get('suggested_filename', 'image.png')
                label = "Missing"
                
                max_chars = max(10, ew // 8)
                if len(suggestion) > max_chars:
                    suggestion = suggestion[:max_chars-3] + "..."
                
                try:
                    font = get_font_cached(font_path, 12)
                except:
                    font = get_font_cached(None, 12)
                
                bbox = draw.textbbox((0, 0), label, font=font)
                text_w = bbox[2] - bbox[0]
                text_x = x + (ew - text_w) // 2
                text_y = y + eh // 3
                draw.text((text_x, text_y), label, fill=(100, 100, 100), font=font)
                
                bbox2 = draw.textbbox((0, 0), suggestion, font=font)
                text_w2 = bbox2[2] - bbox2[0]
                text_x2 = x + (ew - text_w2) // 2
                draw.text((text_x2, text_y + 15), suggestion, fill=(150, 50, 50), font=font)
        
        elif el['type'] == 'text':
            text = user_data.get(el['id'], el.get('text_default', ''))
            
            text_props = el.get('text_props', {}).copy()
            if text_props and text != el.get('text_default', ''):
                if text_props.get('paragraphs'):
                    text_props['paragraphs'][0]['text'] = text
            
            if text_props and text_props.get('paragraphs'):
                text_array = render_text_paragraphs(text_props, ew, eh, font_path)
                
                if text_array.shape[0] > 0 and text_array.shape[1] > 0:
                    y_end = min(y + eh, h)
                    x_end = min(x + ew, w)
                    h_actual = y_end - y
                    w_actual = x_end - x
                    
                    if h_actual > 0 and w_actual > 0:
                        text_crop = text_array[:h_actual, :w_actual]
                        
                        alpha = text_crop[:, :, 3:4].astype(np.float32) / 255.0
                        rgb_layer = text_crop[:, :, :3].astype(np.float32)
                        roi = frame[y:y_end, x:x_end].astype(np.float32)
                        
                        blended = (rgb_layer * alpha + roi * (1 - alpha)).astype(np.uint8)
                        frame[y:y_end, x:x_end] = blended
    
    return frame

def encode_video(frames, fps, output_path):
    h, w = frames[0].shape[:2]
    for codec in ['mp4v', 'avc1']:
        try:
            fourcc = cv2.VideoWriter_fourcc(*codec)
            writer = cv2.VideoWriter(output_path, fourcc, fps, (w, h))
            if writer.isOpened():
                for frame in frames:
                    writer.write(cv2.cvtColor(frame, cv2.COLOR_RGB2BGR))
                writer.release()
                if os.path.exists(output_path) and os.path.getsize(output_path) > 1024:
                    return True
        except:
            continue
    return False

def main():
    st.title("🏭 PPTX Video Factory")
    st.caption("Auto-match → Manual → Upload | Text always renders")
    
    defaults = {
        'layout': None,
        'user_data': {},
        'font_path': find_font(),
        'bg_settings': {'type': 'pptx', 'value': None},
        'asset_sources': {},
        'manual_mappings': {},
        'pptx_path': None,
        'pptx_name': None,
        'templates_dir': "templates",
        'templates': {},
        'current_template': None,
        'selected_element': None,
        'uploaded_images': {},
    }
    for key, val in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = val
    
    with st.sidebar:
        st.header("1. Format")
        layout_name = st.selectbox("Size", list(LAYOUTS.keys()))
        target = LAYOUTS[layout_name]
        
        st.divider()
        st.header("2. Templates Directory")
        templates_dir = st.text_input("Templates folder", value=st.session_state.templates_dir)
        st.session_state.templates_dir = templates_dir
        
        st.divider()
        st.header("3. Upload PPTX")
        uploaded = st.file_uploader("PPTX file", type=['pptx'])
        
        if uploaded and st.button("Extract & Auto-Match"):
            pptx_name = Path(uploaded.name).stem
            st.session_state.pptx_name = pptx_name
            
            extract_dir = Path(tempfile.gettempdir()) / f"pptx_{hashlib.md5(uploaded.getvalue()).hexdigest()[:8]}"
            extract_dir.mkdir(exist_ok=True)
            
            tmp_path = extract_dir / "presentation.pptx"
            with open(tmp_path, 'wb') as f:
                f.write(uploaded.getvalue())
            
            st.session_state.pptx_path = str(tmp_path)
            
            extracted_media = extract_media_from_pptx(str(tmp_path), extract_dir)
            
            auto_discovered = None
            if templates_dir:
                expected_media_path = Path(templates_dir) / pptx_name / "ppt" / "media"
                if expected_media_path.exists():
                    auto_discovered = str(expected_media_path)
                    st.success(f"✓ Auto-discovered: {expected_media_path}")
                else:
                    st.info(f"ℹ️ No auto-folder at {expected_media_path}")
            
            st.session_state.asset_sources = {
                'folder': auto_discovered,
                'extracted_media': extracted_media,
                'uploaded': {}
            }
            
            layout = harvest_ppt(
                str(tmp_path), 
                target, 
                st.session_state.asset_sources,
                st.session_state.manual_mappings
            )
            
            if layout:
                st.session_state.layout = layout
                st.session_state.user_data = {
                    el['id']: el.get('text_default', '')
                    for el in layout['elements']
                    if el.get('type') == 'text'
                }
                st.session_state.bg_settings = {'type': 'pptx', 'value': None}
                
                missing = [el for el in layout['elements'] 
                          if el.get('type') == 'image' and not el.get('image_path')]
                
                matched = sum(1 for e in layout['elements'] if e.get('image_path'))
                total_img = sum(1 for e in layout['elements'] if e.get('type') == 'image')
                
                if missing:
                    st.warning(f"{len(missing)} images missing - upload below")
                else:
                    st.success(f"All {total_img} images matched!")
                
                st.rerun()
        
        if st.session_state.layout:
            layout = st.session_state.layout
            
            st.divider()
            st.header("📊 Stats")
            st.metric("Canvas", f"{layout['canvas']['w']}x{layout['canvas']['h']}")
            
            shapes = sum(1 for e in layout['elements'] if e['type'] == 'shape')
            texts = sum(1 for e in layout['elements'] if e['type'] == 'text')
            images = sum(1 for e in layout['elements'] if e['type'] == 'image')
            matched = sum(1 for e in layout['elements'] if e.get('image_path'))
            
            st.metric("Shapes", shapes)
            st.metric("Text", texts)
            st.metric(f"Images {matched}/{images}", images)
            
            missing = [el for el in layout['elements'] 
                      if el.get('type') == 'image' and not el.get('image_path')]
            
            if missing:
                st.divider()
                st.header(f"⚠️ Missing Images ({len(missing)})")
                st.caption("Auto-match failed. Select existing or upload new.")
                
                manual_pool = {}
                if st.session_state.asset_sources.get('folder'):
                    manual_pool.update(get_images_from_folder(st.session_state.asset_sources['folder']))
                if st.session_state.asset_sources.get('extracted_media'):
                    manual_pool.update(get_images_from_folder(st.session_state.asset_sources['extracted_media']))
                manual_pool.update(st.session_state.uploaded_images)
                
                for el in missing:
                    el_id = el['id']
                    suggestion = el.get('suggested_filename', 'image.png')
                    
                    with st.expander(f"{el_id} (need: {suggestion})", expanded=True):
                        if manual_pool:
                            options = ["(Select existing...)"] + sorted(manual_pool.keys())
                            selected = st.selectbox(
                                "Use existing",
                                options,
                                key=f"manual_select_{el_id}"
                            )
                            
                            if selected != "(Select existing...)":
                                st.session_state.manual_mappings[el_id] = manual_pool[selected]
                                layout = harvest_ppt(
                                    st.session_state.pptx_path,
                                    target,
                                    st.session_state.asset_sources,
                                    st.session_state.manual_mappings
                                )
                                st.session_state.layout = layout
                                st.success(f"Mapped {el_id} to {selected}")
                                st.rerun()
                        
                        st.caption("Or upload new file:")
                        uploaded_file = st.file_uploader(
                            f"Upload for {el_id}",
                            type=['png', 'jpg', 'jpeg', 'gif', 'webp'],
                            key=f"upload_{el_id}"
                        )
                        
                        if uploaded_file:
                            temp_path = Path(tempfile.gettempdir()) / uploaded_file.name
                            with open(temp_path, 'wb') as f:
                                f.write(uploaded_file.getvalue())
                            
                            st.session_state.uploaded_images[uploaded_file.name] = str(temp_path)
                            st.session_state.asset_sources['uploaded'] = st.session_state.uploaded_images
                            st.session_state.manual_mappings[el_id] = str(temp_path)
                            
                            layout = harvest_ppt(
                                st.session_state.pptx_path,
                                target,
                                st.session_state.asset_sources,
                                st.session_state.manual_mappings
                            )
                            st.session_state.layout = layout
                            st.success(f"Uploaded and mapped to {el_id}")
                            st.rerun()
            
            st.divider()
            st.header("📤 General Upload")
            st.caption("Upload images to add to pool")
            
            general_uploads = st.file_uploader(
                "Upload images",
                type=['png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp'],
                accept_multiple_files=True,
                key="general_upload"
            )
            
            if general_uploads:
                added = 0
                for img_file in general_uploads:
                    if img_file.name not in st.session_state.uploaded_images:
                        temp_path = Path(tempfile.gettempdir()) / img_file.name
                        with open(temp_path, 'wb') as f:
                            f.write(img_file.getvalue())
                        st.session_state.uploaded_images[img_file.name] = str(temp_path)
                        added += 1
                
                if added > 0:
                    st.session_state.asset_sources['uploaded'] = st.session_state.uploaded_images
                    st.success(f"Added {added} images to pool")
                    
                    layout = harvest_ppt(
                        st.session_state.pptx_path,
                        target,
                        st.session_state.asset_sources,
                        st.session_state.manual_mappings
                    )
                    st.session_state.layout = layout
                    st.rerun()
            
            st.divider()
            st.header("🎨 Templates")
            
            template_names = list(st.session_state.templates.keys())
            selected_template = st.selectbox(
                "Load Template",
                ["(Current)"] + template_names,
                key="template_selector"
            )
            
            if selected_template != "(Current)" and selected_template in st.session_state.templates:
                if st.button("Load"):
                    template_data = st.session_state.templates[selected_template]
                    st.session_state.manual_mappings = template_data.get('mappings', {})
                    st.session_state.bg_settings = template_data.get('bg', {'type': 'pptx', 'value': None})
                    layout = harvest_ppt(
                        st.session_state.pptx_path,
                        target,
                        st.session_state.asset_sources,
                        st.session_state.manual_mappings
                    )
                    st.session_state.layout = layout
                    st.session_state.current_template = selected_template
                    st.success(f"Loaded: {selected_template}")
                    st.rerun()
            
            new_template_name = st.text_input("Save As", value=st.session_state.current_template or st.session_state.pptx_name or "")
            if st.button("Save Current") and new_template_name:
                st.session_state.templates[new_template_name] = {
                    'mappings': st.session_state.manual_mappings.copy(),
                    'bg': st.session_state.bg_settings.copy(),
                    'pptx_name': st.session_state.pptx_name,
                    'asset_sources': st.session_state.asset_sources
                }
                st.session_state.current_template = new_template_name
                st.success(f"Saved: {new_template_name}")
                st.rerun()
            
            st.divider()
            st.header("🖼️ Background")
            
            bg_opt = st.radio("BG", ["PPTX color", "Solid", "Image"], key="bg_type")
            
            if bg_opt == "PPTX color":
                st.session_state.bg_settings = {'type': 'pptx', 'value': None}
                st.info(f"{layout['pptx_background']}")
            
            elif bg_opt == "Solid":
                hex_c = st.color_picker("Color", "#F5F5F5")
                rgb = tuple(int(hex_c[i:i+2], 16) for i in (1, 3, 5))
                st.session_state.bg_settings = {'type': 'color', 'value': rgb}
            
            elif bg_opt == "Image":
                all_bg = {}
                if st.session_state.asset_sources.get('folder'):
                    all_bg.update(get_images_from_folder(st.session_state.asset_sources['folder']))
                if st.session_state.asset_sources.get('extracted_media'):
                    all_bg.update(get_images_from_folder(st.session_state.asset_sources['extracted_media']))
                all_bg.update(st.session_state.uploaded_images)
                
                if all_bg:
                    bg_name = st.selectbox("Select", list(all_bg.keys()), key="bg_select")
                    if bg_name:
                        try:
                            bg_img = Image.open(all_bg[bg_name]).convert('RGB')
                            st.session_state.bg_settings = {'type': 'image', 'value': bg_img}
                        except Exception as e:
                            st.error(f"Error: {e}")
            
            st.divider()
            st.header("5. Export")
            fps = st.slider("FPS", 6, 60, 30)
            duration = st.slider("Sec", 1, 30, 5)
            st.session_state.export_fps = fps
            st.session_state.export_duration = duration
    
    if not st.session_state.layout:
        st.info("""
        **Workflow:**
        1. Upload PPTX
        2. Tool auto-matches: folder → extracted PPTX media
        3. Missing images shown with suggestions
        4. Manual select existing or upload new
        5. Text always renders
        6. Export
        """)
        return
    
    layout = st.session_state.layout
    
    col_preview, col_controls = st.columns([2, 1])
    
    with col_controls:
        st.subheader("📝 Edit Text")
        for el in layout['elements']:
            if el.get('type') == 'text':
                cur = st.session_state.user_data.get(el['id'], el.get('text_default', ''))
                new = st.text_area(f"{el['id'][:15]}", cur, key=f"t_{el['id']}", height=40)
                st.session_state.user_data[el['id']] = new
        
        if st.button("🔄 Refresh", use_container_width=True, type="primary"):
            st.rerun()
    
    with col_preview:
        bg_info = st.session_state.bg_settings['type']
        sources = [k for k, v in st.session_state.asset_sources.items() if v]
        missing_count = len([e for e in layout['elements'] if e.get('type') == 'image' and not e.get('image_path')])
        st.caption(f"Sources: {', '.join(sources)} | BG: {bg_info} | Missing: {missing_count}")
        
        frame = render_frame(
            layout, 
            st.session_state.user_data, 
            st.session_state.font_path, 
            st.session_state.bg_settings
        )
        st.image(frame, use_container_width=True)
    
    with st.expander(f"Elements ({len(layout['elements'])})"):
        for el in layout['elements']:
            if el['type'] == 'image':
                path = el.get('image_path', 'MISSING')
                status = "✓" if path != 'MISSING' else "✗"
                src = f"({el.get('match_source', '?')})" if path != 'MISSING' else f"[need: {el.get('suggested_filename', 'image.png')}]"
                alt = f" [alt:{el.get('alt_text', 'none')}]" if el.get('alt_text') else ""
                st.text(f"{status} 🖼️ {el['id']}{alt} {src}")
            elif el['type'] == 'text':
                align_info = ""
                if el.get('text_props'):
                    v_align = el['text_props'].get('vertical_align', 'top')
                    if el['text_props'].get('paragraphs'):
                        h_align = el['text_props']['paragraphs'][0].get('align', 'left')
                        align_info = f" [{h_align},{v_align}]"
                st.text(f"📝 {el['id']}{align_info}")
            else:
                st.text(f"⬜ {el['id']}")
    
    st.subheader("Export")
    a, b = st.columns(2)
    
    with a:
        if st.button("🎬 MP4", use_container_width=True, type="primary"):
            with st.spinner(f"Rendering {st.session_state.export_duration}s..."):
                fps = st.session_state.export_fps
                duration = st.session_state.export_duration
                total_frames = int(fps * duration)
                
                frames = []
                for i in range(total_frames):
                    frame = render_frame(
                        layout, 
                        st.session_state.user_data, 
                        st.session_state.font_path, 
                        st.session_state.bg_settings
                    )
                    frames.append(frame)
                
                out = tempfile.mktemp(suffix=".mp4")
                if encode_video(frames, fps, out):
                    with open(out, "rb") as f:
                        st.video(f.read())
                        st.download_button("DL MP4", f, "vid.mp4", mime="video/mp4")
                    os.unlink(out)
    
    with b:
        if st.button("🖼️ PNG", use_container_width=True):
            frame = render_frame(
                layout, 
                st.session_state.user_data, 
                st.session_state.font_path, 
                st.session_state.bg_settings
            )
            img = Image.fromarray(frame)
            
            buf = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            img.save(buf.name)
            with open(buf.name, "rb") as f:
                st.download_button("DL PNG", f, "frame.png", mime="image/png")
            os.unlink(buf.name)

if __name__ == "__main__":
    main()
